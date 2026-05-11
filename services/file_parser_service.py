"""上传资料解析服务。

系统先把用户上传的 txt/md/pdf/docx/pptx 解析为纯文本，再把文本切分写入
Chroma 向量库。解析失败时抛出带文件名和原因的 ValueError，接口层会把原因
保存到 uploaded_files.parse_error，前端可直接展示。
"""

from pathlib import Path

from pptx import Presentation


class FileParserService:
    """把不同格式的文件解析为文本。"""

    SUPPORTED_TYPES = {"txt", "md", "pdf", "docx", "pptx"}

    def parse_file(self, file_path: Path, original_filename: str | None = None) -> str:
        """根据后缀选择解析器，返回清洗后的文本。"""

        filename = original_filename or file_path.name
        suffix = file_path.suffix.lower().lstrip(".")
        if suffix not in self.SUPPORTED_TYPES:
            raise ValueError(f"{filename}：暂不支持 .{suffix or 'unknown'} 格式")

        try:
            if suffix in {"txt", "md"}:
                text = self._parse_text(file_path)
            elif suffix == "pdf":
                text = self._parse_pdf(file_path)
            elif suffix == "docx":
                text = self._parse_docx(file_path)
            elif suffix == "pptx":
                text = self._parse_pptx(file_path)
            else:
                text = ""
        except ValueError:
            raise
        except Exception as exc:
            raise ValueError(f"{filename}：解析失败，原因：{exc}") from exc

        text = self._normalize_text(text)
        if not text:
            raise ValueError(f"{filename}：未提取到有效文本内容")
        return text

    @staticmethod
    def _parse_text(file_path: Path) -> str:
        """txt/md 直接读取文本，兼容常见编码。"""

        for encoding in ("utf-8", "utf-8-sig", "gbk"):
            try:
                return file_path.read_text(encoding=encoding)
            except UnicodeDecodeError:
                continue
        return file_path.read_text(encoding="utf-8", errors="ignore")

    @staticmethod
    def _parse_pdf(file_path: Path) -> str:
        """PDF 优先使用 pypdf 提取文本。"""

        try:
            from pypdf import PdfReader
        except ImportError as exc:
            raise ValueError("缺少 pypdf 依赖，请先执行 python -m pip install -r requirements.txt") from exc

        reader = PdfReader(str(file_path))
        page_texts = []
        for index, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                page_texts.append(f"第 {index} 页\n{text}")
        return "\n\n".join(page_texts)

    @staticmethod
    def _parse_docx(file_path: Path) -> str:
        """docx 使用 python-docx 提取段落和表格文本。"""

        try:
            from docx import Document
        except ImportError as exc:
            raise ValueError("缺少 python-docx 依赖，请先执行 python -m pip install -r requirements.txt") from exc

        document = Document(str(file_path))
        lines = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    lines.append(row_text)
        return "\n".join(lines)

    @staticmethod
    def _parse_pptx(file_path: Path) -> str:
        """pptx 使用 python-pptx 提取每页文本。"""

        presentation = Presentation(str(file_path))
        slide_texts = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slide_texts.append(f"第 {index} 页\n" + "\n".join(texts))
        return "\n\n".join(slide_texts)

    @staticmethod
    def _normalize_text(text: str) -> str:
        """压缩多余空行，减少向量库中的噪声。"""

        lines = [line.strip() for line in (text or "").splitlines()]
        cleaned = []
        previous_blank = False
        for line in lines:
            if not line:
                if not previous_blank:
                    cleaned.append("")
                previous_blank = True
                continue
            cleaned.append(line)
            previous_blank = False
        return "\n".join(cleaned).strip()
