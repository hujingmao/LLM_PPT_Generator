"""python-pptx 生成服务。

PPTService 接收结构化 PPTPlan，根据每页 layout_type 渲染不同版式。
图片接口失败不会影响主流程，系统会自动使用占位图形。
"""

from datetime import datetime
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from config.settings import DEFAULT_TEMPLATE_PATH, OUTPUT_DIR, TEMPLATE_DIR
from models.ppt_schema import PPTPage, PPTPlan
from services.image_service import ImageSearchService
from utils.filename_utils import build_output_filename
from utils.logger import get_logger


logger = get_logger(__name__)


class PPTService:
    """把 PPTPlan 渲染成真实 .pptx 文件。"""

    font_name = "Microsoft YaHei"
    dark = RGBColor(15, 23, 42)
    ink = RGBColor(31, 41, 55)
    muted = RGBColor(100, 116, 139)
    light = RGBColor(248, 250, 252)
    white = RGBColor(255, 255, 255)
    blue = RGBColor(37, 99, 235)
    teal = RGBColor(20, 184, 166)
    coral = RGBColor(244, 114, 82)
    amber = RGBColor(245, 158, 11)

    def __init__(
        self,
        output_dir: Path | None = None,
        template_path: Path | str | None = None,
        image_service: ImageSearchService | None = None,
    ):
        self.output_dir = output_dir or OUTPUT_DIR
        self.output_dir.mkdir(parents=True, exist_ok=True)
        TEMPLATE_DIR.mkdir(parents=True, exist_ok=True)
        self.template_path = Path(template_path) if template_path else (
            DEFAULT_TEMPLATE_PATH if DEFAULT_TEMPLATE_PATH.exists() else None
        )
        self.image_service = image_service

    def generate_ppt(
        self,
        plan: PPTPlan,
        topic: str | None = None,
        template_path: Path | str | None = None,
        use_images: bool = True,
    ) -> Path:
        """根据大纲生成 .pptx，并返回文件路径。"""

        prs, used_template = self._load_presentation(template_path)
        if used_template:
            self._clear_existing_slides(prs)

        pages = plan.pages or [
            PPTPage(
                page_no=1,
                page_title=plan.ppt_title,
                layout_type="cover",
                bullets=[plan.subtitle or plan.theme or "AI generated presentation"],
            )
        ]
        image_service = self.image_service or (ImageSearchService() if use_images else None)
        renderers = {
            "cover": self._add_cover_slide,
            "agenda": self._add_agenda_slide,
            "section": self._add_section_slide,
            "text": self._add_text_slide,
            "image_text": self._add_image_text_slide,
            "three_cards": self._add_three_cards_slide,
            "timeline": self._add_timeline_slide,
            "comparison": self._add_comparison_slide,
            "process": self._add_process_slide,
            "summary": self._add_summary_slide,
            "thanks": self._add_thanks_slide,
        }

        for page in pages:
            image_path = self._fetch_image_for_page(image_service, page, topic or plan.ppt_title) if use_images else None
            renderer = renderers.get(page.layout_type, self._add_text_slide)
            renderer(prs, page, plan, image_path)

        filename = build_output_filename(topic or plan.ppt_title)
        output_path = self.output_dir / filename
        prs.save(str(output_path))
        return output_path

    def _load_presentation(self, template_path: Path | str | None = None) -> tuple[Presentation, Path | None]:
        """加载模板；模板不存在时使用空白宽屏演示文稿。"""

        resolved_template = self._resolve_template_path(template_path or self.template_path)
        if resolved_template:
            return Presentation(str(resolved_template)), resolved_template

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        return prs, None

    @staticmethod
    def _resolve_template_path(template_path: Path | str | None) -> Path | None:
        if not template_path:
            return None
        path = Path(template_path)
        if not path.is_absolute():
            path = TEMPLATE_DIR / path
        path = path.resolve()
        if path.suffix.lower() != ".pptx" or not path.exists():
            return None
        return path

    @staticmethod
    def _clear_existing_slides(prs: Presentation) -> None:
        """清空模板示例页，保留母版和主题资源。"""

        slide_id_list = prs.slides._sldIdLst
        for slide_id in list(slide_id_list):
            prs.part.drop_rel(slide_id.rId)
            slide_id_list.remove(slide_id)

    @staticmethod
    def _blank_layout(prs: Presentation):
        return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    def _add_cover_slide(self, prs: Presentation, page: PPTPage, plan: PPTPlan, image_path: Path | None) -> None:
        """封面页。"""

        slide = prs.slides.add_slide(self._blank_layout(prs))
        self._add_background(slide, prs, self.dark)
        self._add_accent_bar(slide, int(prs.slide_width - Inches(3.6)), Inches(0.9), Inches(2.5), Inches(5.6))

        title = slide.shapes.add_textbox(Inches(0.82), Inches(1.35), int(prs.slide_width * 0.70), Inches(1.6))
        self._set_paragraph(title.text_frame.paragraphs[0], plan.ppt_title or page.page_title, 38, self.white, bold=True)

        subtitle_text = plan.subtitle or " / ".join(page.bullets[:2]) or plan.theme or "AI Presentation"
        subtitle = slide.shapes.add_textbox(Inches(0.86), Inches(3.15), int(prs.slide_width * 0.60), Inches(0.55))
        self._set_paragraph(subtitle.text_frame.paragraphs[0], subtitle_text, 18, RGBColor(203, 213, 225))

        date_box = slide.shapes.add_textbox(Inches(0.86), Inches(6.35), Inches(4.0), Inches(0.35))
        self._set_paragraph(date_box.text_frame.paragraphs[0], datetime.now().strftime("%Y-%m-%d"), 11, RGBColor(148, 163, 184))
        self._add_notes(slide, page)

    def _add_agenda_slide(self, prs: Presentation, page: PPTPage, plan: PPTPlan, image_path: Path | None) -> None:
        """目录页。"""

        slide = prs.slides.add_slide(self._blank_layout(prs))
        self._add_background(slide, prs, self.light)
        self._add_title(slide, "目录", "AGENDA", prs)

        agenda_items = page.bullets or [
            item.page_title for item in plan.pages if item.layout_type not in {"cover", "agenda", "thanks"}
        ]
        for index, item in enumerate(agenda_items[:9], start=1):
            y = Inches(1.45 + (index - 1) * 0.55)
            marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.86), y, Inches(0.38), Inches(0.30))
            marker.fill.solid()
            marker.fill.fore_color.rgb = self.blue if index % 2 else self.teal
            marker.line.fill.background()
            box = slide.shapes.add_textbox(Inches(1.42), y - Inches(0.05), Inches(10.4), Inches(0.42))
            self._set_paragraph(box.text_frame.paragraphs[0], f"{index:02d}  {item}", 18, self.ink)

        self._add_footer(slide, prs, plan.ppt_title)
        self._add_notes(slide, page)

    def _add_section_slide(self, prs: Presentation, page: PPTPage, plan: PPTPlan, image_path: Path | None) -> None:
        """章节过渡页。"""

        slide = prs.slides.add_slide(self._blank_layout(prs))
        self._add_background(slide, prs, self.dark)
        number = slide.shapes.add_textbox(Inches(0.82), Inches(1.05), Inches(1.8), Inches(0.8))
        self._set_paragraph(number.text_frame.paragraphs[0], f"{page.page_no:02d}", 30, self.teal, bold=True)
        title = slide.shapes.add_textbox(Inches(0.82), Inches(2.15), Inches(10.8), Inches(1.0))
        self._set_paragraph(title.text_frame.paragraphs[0], page.page_title, 34, self.white, bold=True)
        self._add_bullets(slide, page.bullets[:3], Inches(0.9), Inches(3.45), Inches(9.8), Inches(1.8), self.white, 18)
        self._add_notes(slide, page)

    def _add_text_slide(self, prs: Presentation, page: PPTPage, plan: PPTPlan, image_path: Path | None) -> None:
        """普通文本页。"""

        slide = prs.slides.add_slide(self._blank_layout(prs))
        self._add_background(slide, prs, self.white)
        self._add_title(slide, page.page_title, "TEXT", prs)
        self._add_bullets(slide, page.bullets, Inches(1.02), Inches(1.55), Inches(11.0), Inches(4.7), self.ink, 19)
        self._add_footer(slide, prs, plan.ppt_title)
        self._add_notes(slide, page)

    def _add_image_text_slide(self, prs: Presentation, page: PPTPage, plan: PPTPlan, image_path: Path | None) -> None:
        """图文页。"""

        slide = prs.slides.add_slide(self._blank_layout(prs))
        self._add_background(slide, prs, self.white)
        self._add_title(slide, page.page_title, "IMAGE + TEXT", prs)
        self._add_bullets(slide, page.bullets[:5], Inches(0.86), Inches(1.55), Inches(5.65), Inches(4.65), self.ink, 17)
        self._add_picture_or_placeholder(slide, image_path, Inches(7.0), Inches(1.35), Inches(5.25), Inches(4.95), page)
        self._add_footer(slide, prs, plan.ppt_title)
        self._add_notes(slide, page)

    def _add_three_cards_slide(self, prs: Presentation, page: PPTPage, plan: PPTPlan, image_path: Path | None) -> None:
        """三栏卡片页。"""

        slide = prs.slides.add_slide(self._blank_layout(prs))
        self._add_background(slide, prs, self.light)
        self._add_title(slide, page.page_title, "THREE CARDS", prs)
        items = (page.bullets or ["核心能力", "实现方式", "应用价值"])[:6]
        columns = [items[0::3], items[1::3], items[2::3]]
        colors = [self.blue, self.teal, self.coral]
        for index, column in enumerate(columns):
            x = Inches(0.82 + index * 4.15)
            card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.55), Inches(3.55), Inches(4.45))
            card.fill.solid()
            card.fill.fore_color.rgb = self.white
            card.line.color.rgb = RGBColor(226, 232, 240)
            label = slide.shapes.add_textbox(x + Inches(0.28), Inches(1.85), Inches(2.9), Inches(0.45))
            self._set_paragraph(label.text_frame.paragraphs[0], f"{index + 1:02d}", 17, colors[index], bold=True)
            self._add_bullets(slide, column, x + Inches(0.28), Inches(2.55), Inches(2.9), Inches(2.8), self.ink, 15)
        self._add_footer(slide, prs, plan.ppt_title)
        self._add_notes(slide, page)

    def _add_timeline_slide(self, prs: Presentation, page: PPTPage, plan: PPTPlan, image_path: Path | None) -> None:
        """时间轴页。"""

        slide = prs.slides.add_slide(self._blank_layout(prs))
        self._add_background(slide, prs, self.white)
        self._add_title(slide, page.page_title, "TIMELINE", prs)
        items = (page.bullets or ["资料解析", "向量入库", "检索增强", "生成导出"])[:5]
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.05), Inches(3.3), Inches(11.1), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(203, 213, 225)
        line.line.fill.background()
        for index, item in enumerate(items):
            x = Inches(1.15 + index * (10.2 / max(len(items) - 1, 1)))
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(3.08), Inches(0.48), Inches(0.48))
            dot.fill.solid()
            dot.fill.fore_color.rgb = [self.blue, self.teal, self.coral, self.amber, self.dark][index % 5]
            dot.line.fill.background()
            box = slide.shapes.add_textbox(x - Inches(0.45), Inches(3.78), Inches(1.8), Inches(1.35))
            self._set_paragraph(box.text_frame.paragraphs[0], item, 14, self.ink, bold=True, alignment=PP_ALIGN.CENTER)
        self._add_footer(slide, prs, plan.ppt_title)
        self._add_notes(slide, page)

    def _add_comparison_slide(self, prs: Presentation, page: PPTPage, plan: PPTPlan, image_path: Path | None) -> None:
        """对比页。"""

        slide = prs.slides.add_slide(self._blank_layout(prs))
        self._add_background(slide, prs, self.light)
        self._add_title(slide, page.page_title, "COMPARISON", prs)
        items = page.bullets or ["传统方式：人工收集资料与排版", "本系统：RAG 增强后自动生成结构化 PPT"]
        mid = max(1, len(items) // 2)
        groups = [items[:mid], items[mid:]]
        titles = ["传统方式", "本系统方案"]
        colors = [self.coral, self.teal]
        for index, group in enumerate(groups):
            x = Inches(0.86 + index * 6.18)
            panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.55), Inches(5.55), Inches(4.55))
            panel.fill.solid()
            panel.fill.fore_color.rgb = self.white
            panel.line.color.rgb = RGBColor(226, 232, 240)
            heading = slide.shapes.add_textbox(x + Inches(0.35), Inches(1.92), Inches(4.8), Inches(0.42))
            self._set_paragraph(heading.text_frame.paragraphs[0], titles[index], 20, colors[index], bold=True)
            self._add_bullets(slide, group, x + Inches(0.35), Inches(2.65), Inches(4.8), Inches(2.7), self.ink, 15)
        self._add_footer(slide, prs, plan.ppt_title)
        self._add_notes(slide, page)

    def _add_process_slide(self, prs: Presentation, page: PPTPage, plan: PPTPlan, image_path: Path | None) -> None:
        """流程页。"""

        slide = prs.slides.add_slide(self._blank_layout(prs))
        self._add_background(slide, prs, self.white)
        self._add_title(slide, page.page_title, "PROCESS", prs)
        items = (page.bullets or ["上传资料", "解析切分", "向量检索", "生成大纲", "导出 PPT"])[:5]
        for index, item in enumerate(items):
            x = Inches(0.78 + index * 2.48)
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.65), Inches(1.95), Inches(1.25))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(239, 246, 255)
            box.line.color.rgb = RGBColor(147, 197, 253)
            text = slide.shapes.add_textbox(x + Inches(0.12), Inches(2.93), Inches(1.72), Inches(0.62))
            self._set_paragraph(text.text_frame.paragraphs[0], item, 14, self.dark, bold=True, alignment=PP_ALIGN.CENTER)
            if index < len(items) - 1:
                arrow = slide.shapes.add_textbox(x + Inches(1.97), Inches(3.03), Inches(0.42), Inches(0.4))
                self._set_paragraph(arrow.text_frame.paragraphs[0], "→", 22, self.teal, bold=True, alignment=PP_ALIGN.CENTER)
        self._add_footer(slide, prs, plan.ppt_title)
        self._add_notes(slide, page)

    def _add_summary_slide(self, prs: Presentation, page: PPTPage, plan: PPTPlan, image_path: Path | None) -> None:
        """总结页。"""

        slide = prs.slides.add_slide(self._blank_layout(prs))
        self._add_background(slide, prs, self.dark)
        title = slide.shapes.add_textbox(Inches(0.82), Inches(0.75), Inches(9.8), Inches(0.8))
        self._set_paragraph(title.text_frame.paragraphs[0], page.page_title or "总结与展望", 31, self.white, bold=True)
        items = (page.bullets or ["完成资料解析、检索增强和 PPT 自动生成闭环", "支持可编辑大纲和多版式导出", "后续可扩展更多模板与协作能力"])[:4]
        for index, item in enumerate(items):
            y = Inches(1.75 + index * 1.05)
            label = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.94), y, Inches(0.58), Inches(0.38))
            label.fill.solid()
            label.fill.fore_color.rgb = [self.blue, self.teal, self.coral, self.amber][index % 4]
            label.line.fill.background()
            box = slide.shapes.add_textbox(Inches(1.78), y - Inches(0.03), Inches(10.2), Inches(0.58))
            self._set_paragraph(box.text_frame.paragraphs[0], item, 18, RGBColor(226, 232, 240))
        self._add_notes(slide, page)

    def _add_thanks_slide(self, prs: Presentation, page: PPTPage, plan: PPTPlan, image_path: Path | None) -> None:
        """致谢页。"""

        slide = prs.slides.add_slide(self._blank_layout(prs))
        self._add_background(slide, prs, self.light)
        title = slide.shapes.add_textbox(Inches(1.2), Inches(2.35), Inches(10.9), Inches(0.92))
        self._set_paragraph(title.text_frame.paragraphs[0], page.page_title or "谢谢聆听", 38, self.dark, bold=True, alignment=PP_ALIGN.CENTER)
        subtitle = slide.shapes.add_textbox(Inches(1.8), Inches(3.45), Inches(9.7), Inches(0.55))
        text = "欢迎老师批评指正" if not page.bullets else " / ".join(page.bullets[:2])
        self._set_paragraph(subtitle.text_frame.paragraphs[0], text, 18, self.muted, alignment=PP_ALIGN.CENTER)
        self._add_notes(slide, page)

    def _fetch_image_for_page(
        self,
        image_service: ImageSearchService | None,
        page: PPTPage,
        topic: str,
    ) -> Path | None:
        """尝试获取配图，失败时返回 None。"""

        if not image_service:
            return None
        fallback_query = f"{topic} {page.page_title}"
        try:
            return image_service.fetch_slide_image(page.keywords, fallback_query=fallback_query)
        except Exception as exc:
            logger.warning("配图获取失败，已降级为占位图形: %s", exc)
            return None

    def _add_background(self, slide, prs: Presentation, color: RGBColor) -> None:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()

    def _add_title(self, slide, title: str, eyebrow: str, prs: Presentation) -> None:
        title_box = slide.shapes.add_textbox(Inches(0.78), Inches(0.55), Inches(8.5), Inches(0.72))
        self._set_paragraph(title_box.text_frame.paragraphs[0], title, 28, self.dark, bold=True)
        tag = slide.shapes.add_textbox(int(prs.slide_width - Inches(2.9)), Inches(0.75), Inches(2.1), Inches(0.35))
        self._set_paragraph(tag.text_frame.paragraphs[0], eyebrow, 10, self.teal, bold=True, alignment=PP_ALIGN.RIGHT)

    def _add_bullets(
        self,
        slide,
        bullets: list[str],
        x,
        y,
        width,
        height,
        color: RGBColor,
        font_size: int,
    ) -> None:
        box = slide.shapes.add_textbox(x, y, width, height)
        text_frame = box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        for index, item in enumerate((bullets or ["待补充内容"])[:8]):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            self._set_paragraph(paragraph, f"• {item}", font_size, color)
            paragraph.space_after = Pt(9)

    def _add_footer(self, slide, prs: Presentation, text: str) -> None:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.72),
            int(prs.slide_height - Inches(0.44)),
            int(prs.slide_width - Inches(1.44)),
            Inches(0.02),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(226, 232, 240)
        line.line.fill.background()
        footer = slide.shapes.add_textbox(Inches(0.72), int(prs.slide_height - Inches(0.35)), Inches(6.2), Inches(0.22))
        self._set_paragraph(footer.text_frame.paragraphs[0], text[:56], 8, self.muted)

    def _add_accent_bar(self, slide, x, y, width, height) -> None:
        colors = [self.blue, self.teal, self.coral, self.amber]
        for index, color in enumerate(colors):
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x + int(width * index / 4),
                y + int(height * index * 0.06),
                int(width / 4),
                int(height * (1 - index * 0.06)),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()

    def _add_picture_or_placeholder(self, slide, image_path: Path | None, x, y, width, height, page: PPTPage) -> None:
        if image_path and image_path.exists():
            self._add_cropped_picture(slide, image_path, x, y, width, height)
            return
        self._add_visual_fallback(slide, x, y, width, height, page)

    def _add_visual_fallback(self, slide, x, y, width, height, page: PPTPage) -> None:
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(239, 246, 255)
        panel.line.color.rgb = RGBColor(191, 219, 254)
        for index, color in enumerate([self.blue, self.teal, self.coral]):
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                int(x + Inches(0.42)),
                int(y + Inches(0.62 + index * 0.62)),
                int(width * (0.72 - index * 0.10)),
                Inches(0.18),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()
        keyword_text = " / ".join(page.keywords[:3]) or page.page_title
        box = slide.shapes.add_textbox(int(x + Inches(0.42)), int(y + height - Inches(1.05)), int(width - Inches(0.84)), Inches(0.5))
        self._set_paragraph(box.text_frame.paragraphs[0], keyword_text[:52], 14, self.dark, bold=True)

    @staticmethod
    def _add_cropped_picture(slide, image_path: Path, x, y, width, height) -> None:
        """按目标框比例裁剪图片，避免拉伸变形。"""

        try:
            with Image.open(image_path) as image:
                image_width, image_height = image.size
            image_ratio = image_width / image_height
            box_ratio = width / height
            picture = slide.shapes.add_picture(str(image_path), x, y, width=width, height=height)
            if image_ratio > box_ratio:
                crop = min((image_ratio - box_ratio) / (2 * image_ratio), 0.49)
                picture.crop_left = crop
                picture.crop_right = crop
            elif image_ratio < box_ratio:
                crop = min((box_ratio - image_ratio) / (2 * box_ratio), 0.49)
                picture.crop_top = crop
                picture.crop_bottom = crop
        except Exception:
            try:
                slide.shapes.add_picture(str(image_path), x, y, width=width, height=height)
            except Exception:
                pass

    def _add_notes(self, slide, page: PPTPage) -> None:
        if page.speaker_notes:
            try:
                slide.notes_slide.notes_text_frame.text = page.speaker_notes
            except Exception:
                logger.debug("写入 speaker notes 失败", exc_info=True)

    def _set_paragraph(
        self,
        paragraph,
        text: str,
        font_size: int,
        color: RGBColor,
        bold: bool = False,
        alignment=None,
    ) -> None:
        paragraph.text = text
        paragraph.font.name = self.font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        if alignment is not None:
            paragraph.alignment = alignment
